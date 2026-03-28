from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # Cyan / teal
ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # Light cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xCA, 0xD9, 0xE8)   # Muted blue-white
GOLD      = RGBColor(0xFF, 0xD6, 0x6E)   # Warm gold
SOFT_BLUE = RGBColor(0x1E, 0x3A, 0x5F)   # Mid navy (boxes)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def bg(slide):
    """Full-slide dark background."""
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)

def accent_bar(slide, top=0.72, height=0.05):
    """Thin horizontal accent bar."""
    add_rect(slide, 0, top, 13.33, height, ACCENT)

def slide_title(slide, title, subtitle=None):
    """Standard slide heading."""
    add_text_box(slide, title, 0.4, 0.15, 12.5, 0.6,
                 font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    accent_bar(slide, top=0.75)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.82, 12.5, 0.4,
                     font_size=14, color=LIGHT_TXT, italic=True)

def bullet_box(slide, bullets, l, t, w, h, start_size=16, color=LIGHT_TXT):
    """Render a list of bullet strings in a text box."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = b
        run.font.size  = Pt(start_size)
        run.font.color.rgb = color

def info_card(slide, l, t, w, h, title, body_lines, title_color=ACCENT):
    """A rounded-looking info card (filled rect + text)."""
    add_rect(slide, l, t, w, h, SOFT_BLUE)
    # Card title
    add_text_box(slide, title, l+0.12, t+0.1, w-0.25, 0.38,
                 font_size=14, bold=True, color=title_color)
    # Card body
    bullet_box(slide, body_lines, l+0.12, t+0.45, w-0.25, h-0.55,
               start_size=12, color=LIGHT_TXT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1  –  Title / Cover
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
bg(sl)

# Big glowing bar on left
add_rect(sl, 0, 0, 0.12, 7.5, ACCENT)

# Main title text
add_text_box(sl, "Multilingual News\nPublishing Portal",
             0.4, 1.3, 12, 1.5, font_size=46, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text_box(sl, "For Independent Reporter",
             0.4, 2.9, 12, 0.55, font_size=26, bold=False,
             color=ACCENT2, align=PP_ALIGN.LEFT, italic=True)

# Accent divider
add_rect(sl, 0.4, 3.55, 6, 0.05, GOLD)

# Meta info
add_text_box(sl, "MCA Semester 4  •  Major Project Presentation",
             0.4, 3.75, 12, 0.45, font_size=16, color=LIGHT_TXT)

# Technology badges
badges = [("React.js", 0.4), ("Node.js", 1.9), ("MongoDB", 3.4),
          ("i18next", 4.9), ("JWT Auth", 6.4)]
for label, lx in badges:
    add_rect(sl, lx, 5.0, 1.3, 0.45, SOFT_BLUE)
    add_text_box(sl, label, lx+0.05, 5.02, 1.2, 0.4,
                 font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(sl, "Slide 1 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2  –  Project Overview
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🌐  Project Overview",
            "A full-stack web application enabling multilingual news publishing")

points = [
    "✦  Single authenticated reporter publishes news in English, Hindi, or Marathi",
    "✦  Auto-translation engine converts articles to all 3 languages automatically",
    "✦  Guest users browse, filter, and read news — no login required",
    "✦  Coverage categories: Local · National · International",
    "✦  Dynamic language switching — readers pick their preferred language in real time",
    "✦  Responsive design: desktop, tablet & mobile",
    "✦  Built as MCA Sem 4 Major Project using modern full-stack technologies",
]
bullet_box(sl, points, 0.5, 1.15, 12.2, 5.8, start_size=17, color=LIGHT_TXT)

add_text_box(sl, "Slide 2 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3  –  Problem Statement
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "❓  Problem Statement")

problems = [
    ("Language Barrier", [
        "• Independent reporters often lack resources",
        "• Writing the same story in 3 languages is",
        "  time-consuming & error-prone",
    ]),
    ("Content Reach", [
        "• Regional readers miss news due to",
        "  language limitations",
        "• No unified portal for tri-lingual content",
    ]),
    ("Manual Effort", [
        "• Translating articles manually takes hours",
        "• No automation → delayed publication",
    ]),
    ("Proposed Solution", [
        "• One-click auto-translation via AI API",
        "• Single platform serving all 3 languages",
        "• Minimal manual effort for the reporter",
    ]),
]

positions = [(0.35, 1.25), (3.55, 1.25), (6.75, 1.25), (9.95, 1.25)]
for (title, lines), (lx, ty) in zip(problems, positions):
    color = GOLD if title == "Proposed Solution" else ACCENT
    info_card(sl, lx, ty, 3.0, 5.7, title, lines, title_color=color)

add_text_box(sl, "Slide 3 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4  –  Tech Stack
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🛠  Technology Stack")

# Two columns
cols = [
    ("⚛  Frontend", [
        "React.js v18  –  UI library",
        "Vite v5       –  Build tool",
        "Tailwind CSS  –  Styling",
        "React Router v6  –  Routing",
        "Axios         –  HTTP client",
        "i18next       –  i18n / translations",
    ]),
    ("🖥  Backend", [
        "Node.js       –  Runtime",
        "Express.js    –  Web framework",
        "MongoDB       –  NoSQL database",
        "Mongoose      –  ODM / schema",
        "JWT           –  Authentication",
        "bcryptjs      –  Password hashing",
        "Multer        –  File uploads",
        "LibreTranslate / Google Translate API",
    ]),
]
for idx, (heading, lines) in enumerate(cols):
    lx = 0.35 + idx * 6.5
    add_rect(sl, lx, 1.2, 6.1, 5.9, SOFT_BLUE)
    add_text_box(sl, heading, lx+0.15, 1.3, 5.7, 0.5,
                 font_size=16, bold=True, color=ACCENT)
    bullet_box(sl, lines, lx+0.15, 1.85, 5.7, 5.0, start_size=14.5, color=LIGHT_TXT)

add_text_box(sl, "Slide 4 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5  –  System Architecture
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🏗  System Architecture")

# Architecture diagram (text-based)
layers = [
    ("CLIENT LAYER", "React.js SPA  •  Vite Dev Server  •  Tailwind CSS  •  i18next", 1.2),
    ("API LAYER",    "Express.js REST API  •  JWT Middleware  •  Multer", 2.7),
    ("SERVICE LAYER","Auto-Translation Engine  •  LibreTranslate / Google Translate API", 4.2),
    ("DATA LAYER",   "MongoDB  •  Mongoose ODM  •  File Storage (uploads/)", 5.7),
]

for (label, detail, ty) in layers:
    add_rect(sl, 1.5, ty, 10.3, 1.1, SOFT_BLUE)
    add_text_box(sl, label, 1.65, ty+0.05, 3, 0.4,
                 font_size=13, bold=True, color=ACCENT)
    add_text_box(sl, detail, 1.65, ty+0.48, 10, 0.5,
                 font_size=13, color=LIGHT_TXT)

# Arrows between layers
for ay in [2.3, 3.8, 5.3]:
    add_rect(sl, 6.55, ay, 0.22, 0.4, ACCENT)

add_text_box(sl, "Slide 5 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6  –  Key Features
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "✨  Key Features")

features = [
    ("🤖  Auto-Translation", [
        "Write once, publish in 3 languages",
        "Powered by Google Translate API",
        "Fallback: LibreTranslate + MyMemory",
        "Re-translates on content edit",
    ]),
    ("🔐  Secure Auth", [
        "JWT tokens (7-day expiry)",
        "bcryptjs password hashing",
        "Protected reporter routes",
        "Single-reporter access model",
    ]),
    ("🗂  News Management", [
        "Create / Edit / Delete articles",
        "Publish / Unpublish toggle",
        "Image upload (5 MB, JPEG/PNG)",
        "Draft & live status tracking",
    ]),
    ("🌍  Multilingual UX", [
        "Dynamic language switcher in Navbar",
        "EN / HI / MR with i18next",
        "Fallback to English if missing",
        "Persistent language preference",
    ]),
    ("📱  Responsive UI", [
        "Tailwind CSS utility-first design",
        "Mobile, tablet & desktop layouts",
        "News card grid layout",
        "Coverage filter (Local/Nat/Int'l)",
    ]),
    ("📊  Reporter Dashboard", [
        "Stats: total / published / draft",
        "One-click publish toggle",
        "Quick edit & delete actions",
        "Manage all articles in one view",
    ]),
]

positions_6 = [
    (0.3, 1.2), (4.55, 1.2), (8.8, 1.2),
    (0.3, 4.0), (4.55, 4.0), (8.8, 4.0),
]
for feat, pos in zip(features, positions_6):
    title, lines = feat
    lx, ty = pos
    info_card(sl, lx, ty, 4.05, 2.5, title, lines)

add_text_box(sl, "Slide 6 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7  –  Database Schema
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🗄  Database Schema  (MongoDB)")

# Reporter Model card
reporter_lines = [
    "username    :  String (required)",
    "email       :  String (unique, required)",
    "password    :  String (hashed via bcrypt)",
    "createdAt   :  Date",
    "updatedAt   :  Date",
]
add_rect(sl, 0.35, 1.2, 5.9, 5.9, SOFT_BLUE)
add_text_box(sl, "👤  Reporter Model", 0.5, 1.3, 5.6, 0.5,
             font_size=16, bold=True, color=ACCENT)
bullet_box(sl, reporter_lines, 0.5, 1.85, 5.6, 2.5, start_size=14, color=LIGHT_TXT)

# News Model card
news_lines = [
    "baseLanguage  :  enum ['en','hi','mr']",
    "title.en / .hi / .mr  :  String",
    "content.en / .hi / .mr  :  String",
    "coverage  :  'local'|'national'|'int'l'",
    "category  :  String",
    "image     :  String (URL / null)",
    "published :  Boolean (default: false)",
    "createdAt / updatedAt  :  Date",
]
add_rect(sl, 6.95, 1.2, 5.9, 5.9, SOFT_BLUE)
add_text_box(sl, "📰  News Model", 7.1, 1.3, 5.6, 0.5,
             font_size=16, bold=True, color=GOLD)
bullet_box(sl, news_lines, 7.1, 1.85, 5.6, 5.0, start_size=14, color=LIGHT_TXT)

add_text_box(sl, "Slide 7 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8  –  API Endpoints
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🔌  RESTful API Endpoints")

auth_endpoints = [
    "POST  /api/auth/register   –  Register reporter (one-time)",
    "POST  /api/auth/login      –  Login & receive JWT",
    "GET   /api/auth/me         –  Get current reporter info  🔒",
]
pub_endpoints = [
    "GET   /api/news            –  All published news (filter: coverage/lang)",
    "GET   /api/news/:id        –  Single news article (public)",
]
priv_endpoints = [
    "GET   /api/news/admin/all  –  All posts incl. drafts  🔒",
    "POST  /api/news            –  Create + auto-translate  🔒",
    "PUT   /api/news/:id        –  Update + re-translate    🔒",
    "DELETE /api/news/:id       –  Delete article           🔒",
    "PATCH  /api/news/:id/publish –  Toggle publish status  🔒",
]

sections = [
    ("🔑  Authentication", auth_endpoints, 1.2),
    ("🌐  Public Endpoints", pub_endpoints, 3.1),
    ("🔒  Protected Endpoints (Reporter)", priv_endpoints, 4.5),
]
for heading, eps, ty in sections:
    add_text_box(sl, heading, 0.4, ty, 12.5, 0.4,
                 font_size=15, bold=True, color=ACCENT)
    bullet_box(sl, eps, 0.5, ty+0.38, 12.3,
               0.45 * len(eps), start_size=13, color=LIGHT_TXT)

add_text_box(sl, "🔒 = JWT token required", 0.4, 6.9, 4, 0.4,
             font_size=11, color=GOLD, italic=True)
add_text_box(sl, "Slide 8 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9  –  Auto-Translation Flow
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
slide_title(sl, "🤖  Auto-Translation Workflow")

steps = [
    ("1", "Reporter writes article in ONE language (EN / HI / MR)"),
    ("2", "Selects base language from dropdown & clicks Publish"),
    ("3", "Backend receives payload → calls Translation API"),
    ("4", "Google Translate API translates title & content"),
    ("5", "Fallback chain: LibreTranslate → MyMemory API"),
    ("6", "All 3 language versions stored in MongoDB"),
    ("7", "Frontend serves content in reader's selected language"),
]

for i, (num, desc) in enumerate(steps):
    ty = 1.2 + i * 0.82
    # Step circle
    add_rect(sl, 0.35, ty, 0.55, 0.55, ACCENT)
    add_text_box(sl, num, 0.35, ty+0.02, 0.55, 0.45,
                 font_size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    # Step text
    add_text_box(sl, desc, 1.05, ty+0.04, 11.5, 0.5,
                 font_size=15, color=LIGHT_TXT)
    # Connector line (except last)
    if i < len(steps)-1:
        add_rect(sl, 0.57, ty+0.55, 0.1, 0.27, ACCENT2)

add_text_box(sl, "Slide 9 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10  –  Conclusion & Future Scope
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
add_rect(sl, 0, 0, 0.12, 7.5, GOLD)
slide_title(sl, "🎯  Conclusion & Future Scope")

conclusions = [
    "✅  Successfully built a trilingual (EN / HI / MR) news portal",
    "✅  Auto-translation eliminates manual effort — reporter writes once",
    "✅  Secure JWT authentication for single-reporter model",
    "✅  RESTful API with full CRUD & image upload support",
    "✅  Responsive, modern UI deployable on Vercel + Railway",
]

future = [
    "🔮  Full-text search across articles",
    "🔮  Rich text editor (draft.js / TipTap)",
    "🔮  Pagination & infinite scroll",
    "🔮  Cloud image storage (Cloudinary)",
    "🔮  Reader comments & social sharing",
    "🔮  Analytics dashboard for reporter",
]

add_rect(sl, 0.35, 1.2, 6.15, 5.7, SOFT_BLUE)
add_text_box(sl, "📌  Achievements", 0.5, 1.3, 5.8, 0.5,
             font_size=16, bold=True, color=ACCENT)
bullet_box(sl, conclusions, 0.5, 1.85, 5.8, 4.8, start_size=14, color=LIGHT_TXT)

add_rect(sl, 6.85, 1.2, 6.1, 5.7, SOFT_BLUE)
add_text_box(sl, "🚀  Future Enhancements", 7.0, 1.3, 5.8, 0.5,
             font_size=16, bold=True, color=GOLD)
bullet_box(sl, future, 7.0, 1.85, 5.8, 4.8, start_size=14, color=LIGHT_TXT)

add_text_box(sl, "Thank You  🙏", 0, 6.6, 13.33, 0.65,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(sl, "Slide 10 / 10", 11.8, 7.0, 1.4, 0.4,
             font_size=10, color=LIGHT_TXT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
out = r"d:\SEM IV\Multilingual_News_Portal_PPT.pptx"
prs.save(out)
print(f"✅  Presentation saved → {out}")
